# Copyright (c) 2025 Bytedance Ltd. and/or its affiliates
# SPDX-License-Identifier: MIT

"""结构化幻灯片数据 -> pptx。纯 python-pptx 直出，不依赖 marp/chrome。

版式一览（data.slides[].layout）：
  agenda / section / bullets / kpi / chart / compare / timeline / table / quote / end
封面由 title/subtitle 自动生成；图表是 pptx 原生图表对象，用户可二次编辑。
图表数据不合法时逐页降级为要点页，绝不让整份 PPT 失败。
"""

import logging
import os
import uuid

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from src.skills.pptx_styles import THEMES, Theme

logger = logging.getLogger(__name__)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_TOP = Inches(1.35)

_CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
}


# ---------- 基础工具 ----------


def _ea_font(run, name: str) -> None:
    """python-pptx 的 font.name 只写 latin 字体，中文要补 <a:ea> 东亚字体。"""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)


def _para(
    tf,
    text: str,
    theme: Theme,
    *,
    size: int = 16,
    bold: bool = False,
    color=None,
    align=PP_ALIGN.LEFT,
    head: bool = False,
    first: bool = False,
):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    f = run.font
    font_name = theme.font_head if head else theme.font_body
    f.name = font_name
    _ea_font(run, font_name)
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color if color is not None else theme.text
    return p


def _textbox(slide, x, y, w, h):
    # 布局算术会产生 float（Emu 相除），python-pptx 只收 int
    box = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    box.text_frame.word_wrap = True
    return box


def _rect(slide, x, y, w, h, color, rounded: bool = False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        int(x), int(y), int(w), int(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------- 页面公共件 ----------


def _content_title(slide, theme: Theme, title: str) -> None:
    box = _textbox(slide, MARGIN, Inches(0.35), SLIDE_W - MARGIN * 2, Inches(0.7))
    _para(box.text_frame, title, theme, size=26, bold=True, color=theme.primary,
          head=True, first=True)
    _rect(slide, MARGIN, Inches(1.08), Inches(0.9), Pt(3.5), theme.accent)


def _insight_bar(slide, theme: Theme, text: str) -> None:
    """页底结论条：一句话洞察，左侧强调色块。"""
    if not text:
        return
    y = SLIDE_H - Inches(0.85)
    _rect(slide, MARGIN, y, SLIDE_W - MARGIN * 2, Inches(0.55), theme.band)
    _rect(slide, MARGIN, y, Inches(0.08), Inches(0.55), theme.accent)
    box = _textbox(slide, MARGIN + Inches(0.25), y, SLIDE_W - MARGIN * 2 - Inches(0.4),
                   Inches(0.55))
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(box.text_frame, str(text), theme, size=13, bold=True, first=True)


def _page_num(slide, theme: Theme, n: int) -> None:
    box = _textbox(slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.45),
                   Inches(0.6), Inches(0.3))
    _para(box.text_frame, str(n), theme, size=10, color=theme.muted,
          align=PP_ALIGN.RIGHT, first=True)


# ---------- 版式 ----------


def _slide_cover(prs, theme: Theme, title: str, subtitle: str, meta: str) -> None:
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.primary)
    _rect(slide, MARGIN, Inches(3.05), Inches(1.2), Pt(4), theme.accent)
    box = _textbox(slide, MARGIN, Inches(3.3), SLIDE_W - MARGIN * 2, Inches(1.6))
    _para(box.text_frame, title, theme, size=40, bold=True, color=theme.on_primary,
          head=True, first=True)
    if subtitle:
        box2 = _textbox(slide, MARGIN, Inches(4.7), SLIDE_W - MARGIN * 2, Inches(0.8))
        _para(box2.text_frame, subtitle, theme, size=18, color=theme.on_primary,
              first=True)
    if meta:
        box3 = _textbox(slide, MARGIN, SLIDE_H - Inches(1.0),
                        SLIDE_W - MARGIN * 2, Inches(0.5))
        _para(box3.text_frame, meta, theme, size=13, color=theme.on_primary, first=True)


def _slide_agenda(prs, theme: Theme, s: dict) -> None:
    slide = _blank(prs)
    _content_title(slide, theme, s.get("title") or "目录")
    points = [str(p) for p in s.get("points", [])][:8]
    box = _textbox(slide, Inches(1.2), CONTENT_TOP + Inches(0.3),
                   SLIDE_W - Inches(2.4), SLIDE_H - CONTENT_TOP - Inches(0.8))
    for i, point in enumerate(points):
        p = _para(box.text_frame, f"{i + 1:02d}　{point}", theme, size=20,
                  bold=True, first=(i == 0))
        p.space_after = Pt(18)


def _slide_section(prs, theme: Theme, s: dict) -> None:
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.primary)
    num = str(s.get("number") or "")
    if num:
        box0 = _textbox(slide, MARGIN, Inches(2.2), Inches(3), Inches(1.0))
        _para(box0.text_frame, num, theme, size=44, bold=True, color=theme.accent,
              head=True, first=True)
    box = _textbox(slide, MARGIN, Inches(3.2), SLIDE_W - MARGIN * 2, Inches(1.2))
    _para(box.text_frame, s.get("title", ""), theme, size=34, bold=True,
          color=theme.on_primary, head=True, first=True)
    if s.get("subtitle"):
        box2 = _textbox(slide, MARGIN, Inches(4.35), SLIDE_W - MARGIN * 2, Inches(0.8))
        _para(box2.text_frame, s["subtitle"], theme, size=16,
              color=theme.on_primary, first=True)


def _slide_bullets(prs, theme: Theme, s: dict) -> None:
    slide = _blank(prs)
    _content_title(slide, theme, s.get("title", ""))
    points = [str(p) for p in s.get("points", [])][:7]
    box = _textbox(slide, MARGIN, CONTENT_TOP + Inches(0.15),
                   SLIDE_W - MARGIN * 2, SLIDE_H - CONTENT_TOP - Inches(1.1))
    for i, point in enumerate(points):
        p = _para(box.text_frame, f"▪  {point}", theme, size=17, first=(i == 0))
        p.space_after = Pt(14)
    _insight_bar(slide, theme, s.get("insight", ""))


def _slide_kpi(prs, theme: Theme, s: dict) -> None:
    slide = _blank(prs)
    _content_title(slide, theme, s.get("title", ""))
    items = s.get("items", [])[:4]
    if not items:
        return
    n = len(items)
    gap = Inches(0.35)
    card_w = (SLIDE_W - MARGIN * 2 - gap * (n - 1)) / n
    card_h = Inches(2.6)
    y = Inches(2.2)
    for i, item in enumerate(items):
        x = MARGIN + (card_w + gap) * i
        _rect(slide, x, y, card_w, card_h, theme.band, rounded=True)
        _rect(slide, x + Inches(0.25), y + Inches(0.3), Inches(0.5), Pt(3), theme.accent)
        box = _textbox(slide, x + Inches(0.2), y + Inches(0.55),
                       card_w - Inches(0.4), Inches(1.1))
        _para(box.text_frame, str(item.get("value", "")), theme, size=34, bold=True,
              color=theme.primary, head=True, first=True)
        box2 = _textbox(slide, x + Inches(0.2), y + Inches(1.7),
                        card_w - Inches(0.4), card_h - Inches(1.8))
        _para(box2.text_frame, str(item.get("label", "")), theme, size=13,
              color=theme.muted, first=True)
    _insight_bar(slide, theme, s.get("insight", ""))


def _slide_chart(prs, theme: Theme, s: dict) -> None:
    slide = _blank(prs)
    _content_title(slide, theme, s.get("title", ""))
    categories = [str(c) for c in s.get("categories", [])]
    series_list = s.get("series", [])
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for ser in series_list:
        values = [float(v) for v in ser.get("values", [])]
        if len(values) != len(categories):
            raise ValueError("series length mismatch")
        chart_data.add_series(str(ser.get("name", "")), values)
    if not categories or not series_list:
        raise ValueError("empty chart data")
    ctype = _CHART_TYPES.get(str(s.get("chart_type", "bar")).lower(),
                             XL_CHART_TYPE.COLUMN_CLUSTERED)
    gf = slide.shapes.add_chart(
        ctype, Inches(1.6), CONTENT_TOP + Inches(0.1),
        SLIDE_W - Inches(3.2), SLIDE_H - CONTENT_TOP - Inches(1.15), chart_data
    )
    chart = gf.chart
    chart.font.size = Pt(12)
    chart.font.name = theme.font_body
    is_pie = ctype == XL_CHART_TYPE.PIE
    chart.has_legend = is_pie or len(series_list) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    if is_pie:
        # 饼图按数据点着色
        points = chart.plots[0].series[0].points
        for i, point in enumerate(points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = theme.palette[i % len(theme.palette)]
        chart.plots[0].has_data_labels = True
        chart.plots[0].data_labels.show_percentage = True
        chart.plots[0].data_labels.number_format_is_linked = False
    else:
        for i, ser in enumerate(chart.series):
            color = theme.palette[i % len(theme.palette)]
            if ctype == XL_CHART_TYPE.LINE_MARKERS:
                # 折线系列填充色会变成面积填充，只染线条
                ser.format.line.color.rgb = color
            else:
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = color
    _insight_bar(slide, theme, s.get("insight", ""))


def _slide_compare(prs, theme: Theme, s: dict) -> None:
    slide = _blank(prs)
    _content_title(slide, theme, s.get("title", ""))
    halves = [s.get("left", {}), s.get("right", {})]
    card_w = (SLIDE_W - MARGIN * 2 - Inches(0.4)) / 2
    y = CONTENT_TOP + Inches(0.1)
    card_h = SLIDE_H - y - Inches(1.05)
    for i, half in enumerate(halves):
        x = MARGIN + (card_w + Inches(0.4)) * i
        _rect(slide, x, y, card_w, card_h, theme.band, rounded=True)
        _rect(slide, x, y, card_w, Inches(0.55),
              theme.accent if i == 0 else theme.primary, rounded=True)
        head = _textbox(slide, x + Inches(0.25), y + Inches(0.06),
                        card_w - Inches(0.5), Inches(0.45))
        _para(head.text_frame, str(half.get("title", "")), theme, size=16, bold=True,
              color=theme.on_primary, head=True, first=True)
        body = _textbox(slide, x + Inches(0.25), y + Inches(0.75),
                        card_w - Inches(0.5), card_h - Inches(0.9))
        for j, point in enumerate([str(p) for p in half.get("points", [])][:6]):
            p = _para(body.text_frame, f"▪  {point}", theme, size=14, first=(j == 0))
            p.space_after = Pt(10)
    _insight_bar(slide, theme, s.get("insight", ""))


def _slide_timeline(prs, theme: Theme, s: dict) -> None:
    slide = _blank(prs)
    _content_title(slide, theme, s.get("title", ""))
    steps = s.get("steps", [])[:6]
    if not steps:
        return
    n = len(steps)
    line_y = Inches(3.1)
    _rect(slide, MARGIN, line_y, SLIDE_W - MARGIN * 2, Pt(2.5), theme.band)
    seg = (SLIDE_W - MARGIN * 2) / n
    dot = Inches(0.22)
    for i, step in enumerate(steps):
        cx = MARGIN + seg * i + seg / 2
        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, int(cx - dot / 2), int(line_y - dot / 2 + Pt(1.25)),
            dot, dot)
        oval.fill.solid()
        oval.fill.fore_color.rgb = theme.accent
        oval.line.color.rgb = theme.on_primary
        oval.line.width = Pt(2)
        oval.shadow.inherit = False
        label = _textbox(slide, int(cx - seg / 2 + Inches(0.1)), Inches(2.35),
                         int(seg - Inches(0.2)), Inches(0.5))
        _para(label.text_frame, str(step.get("label", "")), theme, size=14, bold=True,
              color=theme.primary, align=PP_ALIGN.CENTER, first=True)
        desc = _textbox(slide, int(cx - seg / 2 + Inches(0.1)), Inches(3.55),
                        int(seg - Inches(0.2)), Inches(2.4))
        _para(desc.text_frame, str(step.get("desc", "")), theme, size=12,
              color=theme.muted, align=PP_ALIGN.CENTER, first=True)
    _insight_bar(slide, theme, s.get("insight", ""))


def _slide_table(prs, theme: Theme, s: dict) -> None:
    slide = _blank(prs)
    _content_title(slide, theme, s.get("title", ""))
    headers = [str(h) for h in s.get("headers", [])]
    rows = [[str(c) for c in row] for row in s.get("rows", [])][:9]
    if not headers or not rows:
        raise ValueError("empty table")
    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), MARGIN, CONTENT_TOP + Inches(0.1),
        SLIDE_W - MARGIN * 2,
        min(SLIDE_H - CONTENT_TOP - Inches(1.15), Inches(0.5) * (len(rows) + 1)),
    )
    table = shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme.primary
        cell.text_frame.clear()
        _para(cell.text_frame, h, theme, size=13, bold=True,
              color=theme.on_primary, first=True)
    for i, row in enumerate(rows):
        for j in range(len(headers)):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                theme.band if i % 2 else theme.on_primary)
            cell.text_frame.clear()
            _para(cell.text_frame, row[j] if j < len(row) else "", theme,
                  size=12, first=True)
    _insight_bar(slide, theme, s.get("insight", ""))


def _slide_quote(prs, theme: Theme, s: dict) -> None:
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.band)
    _rect(slide, Inches(1.4), Inches(2.5), Inches(0.1), Inches(2.0), theme.accent)
    box = _textbox(slide, Inches(1.9), Inches(2.4), SLIDE_W - Inches(3.8), Inches(2.0))
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(box.text_frame, str(s.get("text", "")), theme, size=24, bold=True,
          color=theme.primary, head=True, first=True)
    if s.get("source"):
        box2 = _textbox(slide, Inches(1.9), Inches(4.6),
                        SLIDE_W - Inches(3.8), Inches(0.6))
        _para(box2.text_frame, f"—— {s['source']}", theme, size=14,
              color=theme.muted, first=True)


def _slide_end(prs, theme: Theme, s: dict) -> None:
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.primary)
    box = _textbox(slide, 0, Inches(3.1), SLIDE_W, Inches(1.0))
    _para(box.text_frame, str(s.get("title") or "感谢聆听"), theme, size=36,
          bold=True, color=theme.on_primary, head=True,
          align=PP_ALIGN.CENTER, first=True)
    if s.get("subtitle"):
        box2 = _textbox(slide, 0, Inches(4.2), SLIDE_W, Inches(0.6))
        _para(box2.text_frame, s["subtitle"], theme, size=15,
              color=theme.on_primary, align=PP_ALIGN.CENTER, first=True)


_LAYOUTS = {
    "agenda": _slide_agenda,
    "section": _slide_section,
    "bullets": _slide_bullets,
    "kpi": _slide_kpi,
    "chart": _slide_chart,
    "compare": _slide_compare,
    "timeline": _slide_timeline,
    "table": _slide_table,
    "quote": _slide_quote,
    "end": _slide_end,
}


def _fallback_bullets(s: dict) -> dict:
    """图表/表格数据不合法时的降级：把数据摊平成要点页。"""
    points = list(s.get("points", []))
    if s.get("categories") and s.get("series"):
        for ser in s["series"]:
            vals = ser.get("values", [])
            pairs = "，".join(
                f"{c}：{v}" for c, v in zip(s["categories"], vals))
            name = ser.get("name", "")
            points.append(f"{name}　{pairs}" if name else pairs)
    if s.get("headers") and s.get("rows"):
        for row in s["rows"][:7]:
            points.append(" / ".join(str(c) for c in row))
    return {"title": s.get("title", ""), "points": points,
            "insight": s.get("insight", "")}


def build_pptx(data: dict, theme_name: str = "business") -> str:
    """deck 数据 -> pptx 文件路径。data: {title, subtitle?, meta?, slides:[...]}."""
    theme = THEMES.get(theme_name, THEMES["business"])
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_cover(prs, theme, data.get("title", "演示文稿"),
                 data.get("subtitle", ""), data.get("meta", ""))
    slides = data.get("slides", [])
    for s in slides:
        layout = str(s.get("layout", "bullets")).lower()
        render = _LAYOUTS.get(layout, _slide_bullets)
        try:
            render(prs, theme, s)
        except Exception:
            logger.warning("slide layout %r failed, fallback to bullets", layout,
                           exc_info=True)
            _slide_bullets(prs, theme, _fallback_bullets(s))
        _page_num(prs.slides[-1], theme, len(prs.slides) - 1)
    if not slides or str(slides[-1].get("layout", "")).lower() != "end":
        _slide_end(prs, theme, {})

    path = os.path.join(os.getcwd(), f"slides_{uuid.uuid4().hex}.pptx")
    prs.save(path)
    return path
