# Copyright (c) 2025 Bytedance Ltd. and/or its affiliates
# SPDX-License-Identifier: MIT

"""PPT 直出管线（pptx_export/pptx_styles）与文档图表嵌入的单元测试."""

import os

from src.skills.pptx_export import build_pptx
from src.skills.pptx_styles import THEMES, resolve_theme

FULL_DECK = {
    "title": "2026 上半年工作汇报",
    "subtitle": "运营中心",
    "meta": "汇报人：张三 · 2026 年 7 月",
    "slides": [
        {"layout": "agenda", "title": "目录", "points": ["目标回顾", "完成情况", "问题分析", "下步计划"]},
        {"layout": "section", "number": "01", "title": "目标回顾", "subtitle": "年初定的三件事"},
        {"layout": "bullets", "title": "半年目标", "points": ["GMV 破亿", "复购率提升至 35%", "新客成本下降两成"], "insight": "三大目标完成两项"},
        {"layout": "kpi", "title": "核心指标", "items": [
            {"value": "1.2 亿", "label": "GMV（同比 +18%）"},
            {"value": "36%", "label": "复购率"},
            {"value": "-12%", "label": "新客成本"},
        ], "insight": "GMV 与复购率均超额达成"},
        {"layout": "chart", "chart_type": "bar", "title": "季度 GMV（万元）",
         "categories": ["Q1", "Q2"], "series": [{"name": "GMV", "values": [5200, 6800]}],
         "insight": "Q2 环比增长 31%"},
        {"layout": "chart", "chart_type": "pie", "title": "渠道占比",
         "categories": ["直播", "搜索", "私域"], "series": [{"name": "占比", "values": [55, 30, 15]}]},
        {"layout": "chart", "chart_type": "line", "title": "月度复购率（%）",
         "categories": ["1月", "2月", "3月"], "series": [
             {"name": "今年", "values": [30, 33, 36]}, {"name": "去年", "values": [28, 29, 30]}]},
        {"layout": "compare", "title": "两种打法对比",
         "left": {"title": "直播投放", "points": ["起量快", "成本高"]},
         "right": {"title": "私域运营", "points": ["成本低", "见效慢"]},
         "insight": "组合使用，直播拉新私域留存"},
        {"layout": "timeline", "title": "下半年节奏", "steps": [
            {"label": "Q3", "desc": "私域体系搭建"}, {"label": "双11", "desc": "全渠道冲刺"},
            {"label": "Q4", "desc": "复盘与沉淀"}]},
        {"layout": "table", "title": "重点项目排期",
         "headers": ["项目", "负责人", "截止"], "rows": [["会员体系", "李四", "8月"], ["积分商城", "王五", "9月"]]},
        {"layout": "quote", "text": "增长的本质是复购", "source": "内部共识"},
        {"layout": "end", "title": "感谢聆听", "subtitle": "欢迎指正"},
    ],
}


class TestPptxExport:
    def test_full_deck_all_layouts(self):
        path = build_pptx(FULL_DECK, theme_name="business")
        assert os.path.exists(path) and os.path.getsize(path) > 30_000
        from pptx import Presentation

        prs = Presentation(path)
        # 封面 + 12 个内容页（end 已含，不再追加）
        assert len(list(prs.slides)) == 13
        # 原生图表对象存在（bar/pie/line 三页）
        chart_slides = [
            s for s in prs.slides if any(sh.has_chart for sh in s.shapes)
        ]
        assert len(chart_slides) == 3
        # 表格页存在
        assert any(
            any(sh.has_table for sh in s.shapes) for s in prs.slides
        )
        os.remove(path)

    def test_bad_chart_falls_back_not_crash(self):
        deck = {
            "title": "T",
            "slides": [
                {"layout": "chart", "chart_type": "bar", "title": "坏图",
                 "categories": ["a", "b"], "series": [{"name": "s", "values": ["x", 2]}],
                 "insight": "还有结论"},
            ],
        }
        path = build_pptx(deck)
        from pptx import Presentation

        prs = Presentation(path)
        # 封面 + 降级要点页 + 自动结尾页
        assert not any(sh.has_chart for s in prs.slides for sh in s.shapes)
        os.remove(path)

    def test_unknown_layout_renders_bullets(self):
        deck = {"title": "T", "slides": [{"layout": "whatever", "title": "x", "points": ["a"]}]}
        path = build_pptx(deck)
        assert os.path.exists(path)
        os.remove(path)

    def test_theme_mapping(self):
        assert resolve_theme(None).name == "business"
        assert resolve_theme("workreport").name == "business"
        assert resolve_theme("pitch").name == "consult"
        assert resolve_theme("courseware").name == "academic"
        assert resolve_theme("unknown-sub").name == "business"
        assert set(THEMES) == {"business", "consult", "academic"}


class TestMckAdapter:
    def test_full_storyline_renders(self):
        from src.skills.mck_adapter import build_mck_deck

        deck = {
            "slides": [
                {"layout": "cover", "title": "复盘", "subtitle": "s", "author": "", "date": ""},
                {"layout": "toc", "items": [{"num": "", "title": "大盘", "desc": "d"}]},
                {"layout": "big_number", "title": "用户破千", "number": "1,200", "unit": "人",
                 "description": "d", "details": ["a", "b"]},
                {"layout": "donut", "title": "占比集中",
                 "segments": [{"label": "PPT", "pct": 45}, {"label": "其他", "pct": 55}],
                 "center_label": "45%", "center_sub": "PPT"},
                {"layout": "closing", "title": "", "message": ""},
            ]
        }
        path = build_mck_deck(deck)
        from pptx import Presentation

        assert len(list(Presentation(path).slides)) == 5
        os.remove(path)

    def test_guard_rails(self):
        from src.skills.mck_adapter import _chevron, _timeline, _title

        # 行动标题超 40 字截断
        assert len(_title({"title": "超" * 60})) == 40
        # chevron 超 5 步截断、label 去换行
        out = _chevron({"title": "t", "steps": [
            {"label": f"第\n{i}", "title": "x", "desc": "y"} for i in range(8)]})
        assert len(out["steps"]) == 5
        assert all("\n" not in s[0] for s in out["steps"])
        # timeline 最后一个 label 压到 6 字内
        out = _timeline({"title": "t", "milestones": [
            {"label": "很长很长的标签啊", "desc": "d"}]})
        assert len(out["milestones"][-1][0]) <= 6

    def test_empty_or_unknown_raises(self):
        import pytest

        from src.skills.mck_adapter import build_mck_deck

        with pytest.raises(ValueError):
            build_mck_deck({"slides": [{"layout": "nope", "title": "x"}]})

    def test_auto_cover_and_closing(self):
        from src.skills.mck_adapter import build_mck_deck

        path = build_mck_deck({"slides": [
            {"layout": "big_number", "title": "只有一页", "number": "1", "unit": "",
             "description": ""}]})
        from pptx import Presentation

        assert len(list(Presentation(path).slides)) == 3  # 自动补 cover + closing
        os.remove(path)


class TestDocxStyling:
    def test_footer_page_number_and_styles(self):
        from src.skills.docx_export import build_sections_docx

        path = build_sections_docx(
            {"title": "周报", "sections": [{"heading": "一、进展", "content": "推进中"}]}
        )
        import docx

        d = docx.Document(path)
        footer_xml = d.sections[0].footer.paragraphs[0]._p.xml
        assert "PAGE" in footer_xml  # 页码域存在
        assert d.styles["Heading 1"].font.name == "微软雅黑"
        os.remove(path)

    def test_notice_variant_red_title_and_indent(self):
        from src.skills.docx_export import build_sections_docx

        path = build_sections_docx(
            {"title": "关于放假的通知", "sections": [{"heading": "一、安排", "content": "正文段落"}]},
            variant="notice",
        )
        import docx
        from docx.shared import RGBColor

        d = docx.Document(path)
        title_run = d.paragraphs[0].runs[0]
        assert title_run.font.color.rgb == RGBColor(0xC0, 0x00, 0x00)
        body = [p for p in d.paragraphs if p.text == "正文段落"][0]
        assert body.paragraph_format.first_line_indent is not None
        os.remove(path)


class TestDocChart:
    def test_sections_docx_with_chart_embeds_image(self):
        from src.skills.charts import _cjk_font
        from src.skills.docx_export import build_sections_docx

        data = {
            "title": "月度经营分析",
            "sections": [
                {"heading": "一、营收", "content": "整体向好",
                 "chart": {"type": "bar", "title": "月度营收（万元）",
                           "categories": ["4月", "5月", "6月"],
                           "series": [{"name": "营收", "values": [120, 135, 160]}]}},
            ],
        }
        path = build_sections_docx(data)
        import docx

        d = docx.Document(path)
        if _cjk_font():
            assert len(d.inline_shapes) == 1  # 图表以内嵌图片存在
        else:
            texts = "\n".join(p.text for p in d.paragraphs)
            assert "4月：120" in texts  # 无 CJK 字体时降级为文字
        os.remove(path)

    def test_chart_bad_data_returns_none(self):
        from src.skills.charts import render_chart_png

        assert render_chart_png({"type": "bar", "categories": [], "series": []}) is None
        assert render_chart_png(
            {"type": "bar", "categories": ["a"], "series": [{"values": [1, 2]}]}
        ) is None
